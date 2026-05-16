#!/usr/bin/env python3
"""Generate a Dry Eye Disease Early Prediction PPT."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

SCRIPT_DIR = os.path.dirname(os.path.abspath(__file__))
PROJECT_DIR = os.path.dirname(SCRIPT_DIR)
OUT = os.path.join(PROJECT_DIR, "Dry_Eye_Early_Prediction.pptx")

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
W = prs.slide_width
H = prs.slide_height

# Colors
DARK = RGBColor(0x0F, 0x17, 0x2A)
BLUE = RGBColor(0x1E, 0x3A, 0x8A)
ACCENT = RGBColor(0x38, 0xBD, 0xF8)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT = RGBColor(0xCB, 0xD5, 0xE1)
ORANGE = RGBColor(0xF9, 0x73, 0x16)
GREEN = RGBColor(0x22, 0xC5, 0x5E)
RED = RGBColor(0xEF, 0x44, 0x44)

def add_bg(slide, color=DARK):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_box(slide, left, top, w, h, color):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    return s

def add_rounded_box(slide, left, top, w, h, color):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    return s

def add_text(slide, left, top, w, h, text, size=18, color=WHITE, bold=False, align=PP_ALIGN.LEFT):
    tx = slide.shapes.add_textbox(left, top, w, h)
    tf = tx.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = align
    return tx

def add_para(tf, text, size=16, color=WHITE, bold=False, align=PP_ALIGN.LEFT, space_before=Pt(6)):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = align
    p.space_before = space_before
    return p

# ─── SLIDE 1: Title ───
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl)
add_box(sl, Inches(0), Inches(0), W, Inches(0.08), ACCENT)
add_text(sl, Inches(1), Inches(1.5), Inches(11), Inches(1.2),
         "Early-Stage Prediction of Dry Eye Disease", 42, WHITE, True, PP_ALIGN.CENTER)
add_text(sl, Inches(1), Inches(2.8), Inches(11), Inches(0.8),
         "A Machine Learning Approach Using Lifestyle, Sleep & Digital Habits", 22, ACCENT, False, PP_ALIGN.CENTER)
add_text(sl, Inches(1), Inches(4.2), Inches(11), Inches(0.6),
         "Problem Statement & Research Proposal", 18, LIGHT, False, PP_ALIGN.CENTER)
add_box(sl, Inches(5.5), Inches(5.2), Inches(2.3), Inches(0.04), ACCENT)
add_text(sl, Inches(1), Inches(5.8), Inches(11), Inches(0.5),
         "2026", 16, LIGHT, False, PP_ALIGN.CENTER)

# ─── SLIDE 2: The Problem ───
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl)
add_box(sl, Inches(0), Inches(0), Inches(0.12), H, ACCENT)
add_text(sl, Inches(0.8), Inches(0.4), Inches(6), Inches(0.7), "The Problem", 32, ACCENT, True)

items = [
    ("344 Million+", "people affected globally by Dry Eye Disease (DED) — WHO 2023"),
    ("75% Undiagnosed", "patients remain undiagnosed until significant corneal damage occurs"),
    ("$3.8 Billion", "annual economic burden in the US alone (direct medical costs)"),
    ("Screen Time Crisis", "avg. 7+ hrs/day digital exposure — #1 modifiable risk factor in <40 age group"),
]
for i, (title, desc) in enumerate(items):
    y = Inches(1.5 + i * 1.3)
    b = add_rounded_box(sl, Inches(0.8), y, Inches(11.5), Inches(1.1), RGBColor(0x1E, 0x29, 0x3B))
    tx = slide.shapes.add_textbox(Inches(1.2), y + Inches(0.15), Inches(10.5), Inches(0.9)) if False else sl.shapes.add_textbox(Inches(1.2), y + Inches(0.15), Inches(10.5), Inches(0.9))
    tf = tx.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(22)
    p.font.color.rgb = ORANGE
    p.font.bold = True
    add_para(tf, desc, 15, LIGHT, space_before=Pt(4))

add_text(sl, Inches(0.8), Inches(6.5), Inches(11), Inches(0.6),
         "⚠  Current diagnosis relies on subjective slit-lamp exams — invasive, expensive, and too late.", 14, RED, True)

# ─── SLIDE 3: Why Early Detection Matters ───
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl)
add_box(sl, Inches(0), Inches(0), Inches(0.12), H, GREEN)
add_text(sl, Inches(0.8), Inches(0.4), Inches(8), Inches(0.7), "Why Early Detection Matters", 32, GREEN, True)

reasons = [
    "Irreversible Damage — Chronic DED leads to corneal scarring, vision impairment & reduced quality of life",
    "Preventable Progression — 80% of early-stage cases are manageable with lifestyle modifications alone",
    "Digital Generation at Risk — Gen Z & Millennials show 2.5× higher DED incidence vs. prior generations",
    "Silent Onset — Symptoms like occasional dryness are dismissed; by the time patients seek help, damage is advanced",
    "Cost Savings — Early intervention costs ~$200/yr vs. $2,000+/yr for late-stage treatment",
    "Comorbidity Signal — DED correlates with sleep disorders, stress, cardiovascular issues — early flag for systemic health",
]
for i, r in enumerate(reasons):
    y = Inches(1.4 + i * 0.9)
    add_text(sl, Inches(1.2), y, Inches(11), Inches(0.8), f"▸  {r}", 16, WHITE if i % 2 == 0 else LIGHT)

# ─── SLIDE 4: Problem Statement (Formal) ───
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl)
add_box(sl, Inches(0), Inches(0), W, Inches(0.08), ORANGE)
add_text(sl, Inches(0.8), Inches(0.5), Inches(11), Inches(0.7), "Problem Statement", 34, ORANGE, True)

stmt = (
    '"Can we predict the onset of Dry Eye Disease at an early stage using '
    'non-invasive, routinely available patient data — including sleep patterns, '
    'stress levels, digital screen habits, lifestyle factors, and ocular symptoms '
    '— thereby enabling timely preventive intervention before irreversible '
    'corneal damage occurs?"'
)
b = add_rounded_box(sl, Inches(0.8), Inches(1.6), Inches(11.7), Inches(2.2), RGBColor(0x1E, 0x29, 0x3B))
add_text(sl, Inches(1.3), Inches(1.8), Inches(10.7), Inches(2.0), stmt, 20, WHITE, True, PP_ALIGN.CENTER)

sub_points = [
    "Objective 1 — Build a classification model achieving ≥ 75% F1-Score on unseen patient data",
    "Objective 2 — Identify top risk factors (feature importance) to guide clinical screening protocols",
    "Objective 3 — Create an interpretable risk-scoring tool usable in primary care & telemedicine settings",
    "Objective 4 — Validate across demographics (age, gender, lifestyle) for generalizability",
]
for i, sp in enumerate(sub_points):
    add_text(sl, Inches(1.2), Inches(4.2 + i * 0.65), Inches(11), Inches(0.55), sp, 16, ACCENT if i == 0 else LIGHT)

# ─── SLIDE 5: Dataset Overview ───
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl)
add_box(sl, Inches(0), Inches(0), Inches(0.12), H, ACCENT)
add_text(sl, Inches(0.8), Inches(0.4), Inches(8), Inches(0.7), "Dataset Overview", 32, ACCENT, True)
add_text(sl, Inches(0.8), Inches(1.1), Inches(10), Inches(0.5),
         "20,000 patient records  ×  26 attributes  |  Binary target: Dry Eye Disease (Y/N)", 16, LIGHT)

cats = [
    ("Demographics", "Gender, Age, Height, Weight", ACCENT),
    ("Sleep & Stress", "Sleep duration, quality, disorder, wake-up, daytime sleepiness, stress level", RGBColor(0xA7, 0x8B, 0xFA)),
    ("Cardiovascular", "Blood pressure (systolic/diastolic), heart rate, daily steps, physical activity", GREEN),
    ("Lifestyle Habits", "Caffeine, alcohol, smoking, medication, medical issues", ORANGE),
    ("Digital Exposure", "Smart device before bed, avg screen time, blue-light filter usage", RED),
    ("Ocular Symptoms", "Eye-strain/discomfort, redness, itchiness/irritation", RGBColor(0xF4, 0x72, 0xB6)),
]
for i, (cat, desc, clr) in enumerate(cats):
    y = Inches(1.8 + i * 0.85)
    add_rounded_box(sl, Inches(0.8), y, Inches(11.5), Inches(0.72), RGBColor(0x1E, 0x29, 0x3B))
    add_text(sl, Inches(1.2), y + Inches(0.08), Inches(3), Inches(0.35), cat, 16, clr, True)
    add_text(sl, Inches(1.2), y + Inches(0.38), Inches(10.5), Inches(0.3), desc, 13, LIGHT)

# ─── SLIDE 6: Proposed Methodology ───
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl)
add_box(sl, Inches(0), Inches(0), Inches(0.12), H, ORANGE)
add_text(sl, Inches(0.8), Inches(0.4), Inches(8), Inches(0.7), "Proposed Methodology", 32, ORANGE, True)

steps = [
    ("01", "Data Collection & Cleaning", "Handle missing values, parse blood pressure, encode categoricals, normalize numerics"),
    ("02", "Exploratory Data Analysis", "Distribution analysis, correlation heatmaps, class imbalance assessment (65:35 split)"),
    ("03", "Feature Engineering", "Statistical feature selection — Cramér's V (categorical), Point-Biserial r (numeric), Chi² tests"),
    ("04", "Model Training", "6 classifiers: Logistic Regression, Decision Tree, Random Forest, Gradient Boosting, XGBoost, KNN"),
    ("05", "Evaluation & Selection", "Metrics: Accuracy, Precision, Recall, F1-Score, AUC-ROC, 5-fold stratified cross-validation"),
    ("06", "Interpretation & Deployment", "Feature importance ranking, confusion matrix analysis, risk-score tool for clinicians"),
]
for i, (num, title, desc) in enumerate(steps):
    y = Inches(1.3 + i * 0.95)
    add_rounded_box(sl, Inches(0.8), y, Inches(11.5), Inches(0.82), RGBColor(0x1E, 0x29, 0x3B))
    add_text(sl, Inches(1.0), y + Inches(0.12), Inches(0.7), Inches(0.55), num, 24, ORANGE, True)
    add_text(sl, Inches(1.8), y + Inches(0.08), Inches(4), Inches(0.35), title, 17, WHITE, True)
    add_text(sl, Inches(1.8), y + Inches(0.42), Inches(10), Inches(0.35), desc, 13, LIGHT)

# ─── SLIDE 7: ML Models ───
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl)
add_box(sl, Inches(0), Inches(0), Inches(0.12), H, ACCENT)
add_text(sl, Inches(0.8), Inches(0.4), Inches(8), Inches(0.7), "Machine Learning Models", 32, ACCENT, True)

models_info = [
    ("Logistic Regression", "Linear", "Baseline; interpretable coefficients", RGBColor(0x3B,0x82,0xF6)),
    ("Decision Tree", "Tree", "Explainable; visual decision paths", GREEN),
    ("Random Forest", "Ensemble", "Bagging; reduces overfitting", RGBColor(0x10,0xB9,0x81)),
    ("Gradient Boosting", "Ensemble", "Sequential error correction", ORANGE),
    ("XGBoost", "Ensemble", "State-of-the-art tabular performance", RED),
    ("K-Nearest Neighbors", "Instance", "Non-parametric; distance-based", RGBColor(0x8B,0x5C,0xF6)),
]
for i, (name, cat, desc, clr) in enumerate(models_info):
    col = i % 3
    row = i // 3
    x = Inches(0.8 + col * 4.1)
    y = Inches(1.3 + row * 2.8)
    add_rounded_box(sl, x, y, Inches(3.8), Inches(2.4), RGBColor(0x1E, 0x29, 0x3B))
    add_box(sl, x, y, Inches(3.8), Inches(0.06), clr)
    add_text(sl, x + Inches(0.3), y + Inches(0.3), Inches(3.2), Inches(0.4), name, 18, WHITE, True)
    add_text(sl, x + Inches(0.3), y + Inches(0.8), Inches(3.2), Inches(0.35), f"Type: {cat}", 14, clr, True)
    add_text(sl, x + Inches(0.3), y + Inches(1.3), Inches(3.2), Inches(0.8), desc, 14, LIGHT)

# ─── SLIDE 8: Expected Outcomes ───
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl)
add_box(sl, Inches(0), Inches(0), Inches(0.12), H, GREEN)
add_text(sl, Inches(0.8), Inches(0.4), Inches(8), Inches(0.7), "Expected Outcomes & Impact", 32, GREEN, True)

outcomes = [
    ("🎯  Accurate Predictor", "A validated ML model achieving ≥75% F1-Score for early DED risk classification"),
    ("📊  Risk Factor Ranking", "Quantified feature importance — identifying which lifestyle habits matter most"),
    ("🏥  Clinical Decision Tool", "A deployable risk-scoring system for ophthalmologists & primary care physicians"),
    ("📱  Telemedicine Ready", "Non-invasive inputs only — can be integrated into mobile health questionnaires"),
    ("💰  Cost Reduction", "Shift from reactive treatment ($2K+/yr) to preventive care ($200/yr)"),
    ("🔬  Research Contribution", "Published dataset analysis & model benchmarks for the DED research community"),
]
for i, (title, desc) in enumerate(outcomes):
    col = i % 2
    row = i // 2
    x = Inches(0.8 + col * 6.2)
    y = Inches(1.3 + row * 1.8)
    add_rounded_box(sl, x, y, Inches(5.8), Inches(1.5), RGBColor(0x1E, 0x29, 0x3B))
    add_text(sl, x + Inches(0.3), y + Inches(0.15), Inches(5.2), Inches(0.45), title, 18, WHITE, True)
    add_text(sl, x + Inches(0.3), y + Inches(0.7), Inches(5.2), Inches(0.7), desc, 14, LIGHT)

# ─── SLIDE 9: Timeline ───
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl)
add_box(sl, Inches(0), Inches(0), Inches(0.12), H, RGBColor(0xA7, 0x8B, 0xFA))
add_text(sl, Inches(0.8), Inches(0.4), Inches(8), Inches(0.7), "Project Timeline", 32, RGBColor(0xA7, 0x8B, 0xFA), True)

phases = [
    ("Week 1-2", "Data Collection & EDA", "Dataset acquisition, cleaning, exploratory analysis, visualization"),
    ("Week 3-4", "Feature Engineering", "Statistical tests, feature selection, correlation analysis"),
    ("Week 5-6", "Model Development", "Train 6 classifiers, hyperparameter tuning, cross-validation"),
    ("Week 7-8", "Evaluation & Analysis", "Performance comparison, confusion matrices, ROC curves"),
    ("Week 9-10", "Tool Development", "Risk-scoring interface, clinical dashboard, documentation"),
    ("Week 11-12", "Report & Presentation", "Final report, research paper draft, presentation preparation"),
]
for i, (wk, title, desc) in enumerate(phases):
    y = Inches(1.3 + i * 0.95)
    add_rounded_box(sl, Inches(0.8), y, Inches(11.5), Inches(0.82), RGBColor(0x1E, 0x29, 0x3B))
    add_text(sl, Inches(1.0), y + Inches(0.12), Inches(1.6), Inches(0.35), wk, 15, ORANGE, True)
    add_text(sl, Inches(2.8), y + Inches(0.08), Inches(4), Inches(0.35), title, 17, WHITE, True)
    add_text(sl, Inches(2.8), y + Inches(0.42), Inches(9), Inches(0.35), desc, 13, LIGHT)

# ─── SLIDE 10: Thank You ───
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl)
add_box(sl, Inches(0), Inches(3.4), W, Inches(0.06), ACCENT)
add_text(sl, Inches(1), Inches(2.0), Inches(11), Inches(1.0),
         "Thank You", 48, WHITE, True, PP_ALIGN.CENTER)
add_text(sl, Inches(1), Inches(3.8), Inches(11), Inches(0.6),
         "Questions & Discussion", 24, ACCENT, False, PP_ALIGN.CENTER)


# Save
prs.save(OUT)
print(f"✅ PPT saved: {OUT}")
