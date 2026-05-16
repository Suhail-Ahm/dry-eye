#!/usr/bin/env python3
"""
Dry Eye Disease — Minimal Modern Presentation with Visual Charts
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os

SCRIPT_DIR = os.path.dirname(os.path.abspath(__file__))
PROJECT_DIR = os.path.dirname(SCRIPT_DIR)
CHART_DIR = os.path.join(SCRIPT_DIR, "chart_images")
OUT = os.path.join(PROJECT_DIR, "Dry_Eye_Minimal.pptx")

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
W = prs.slide_width
H = prs.slide_height

# ── Palette ──────────────────────────────────────────────────────────────
BG      = RGBColor(0x09, 0x09, 0x0B)
SURFACE = RGBColor(0x18, 0x18, 0x1B)
BORDER  = RGBColor(0x27, 0x27, 0x2A)
WHITE   = RGBColor(0xFA, 0xFA, 0xFA)
MUTED   = RGBColor(0xA1, 0xA1, 0xAA)
DIM     = RGBColor(0x71, 0x71, 0x7A)
ACCENT  = RGBColor(0x38, 0xBD, 0xF8)
TEAL    = RGBColor(0x2D, 0xD4, 0xBF)
AMBER   = RGBColor(0xFB, 0xBF, 0x24)
ROSE    = RGBColor(0xFB, 0x71, 0x85)
VIOLET  = RGBColor(0xA7, 0x8B, 0xFA)
GREEN   = RGBColor(0x4A, 0xDE, 0x80)

def bg_fill(slide):
    fill = slide.background.fill; fill.solid(); fill.fore_color.rgb = BG

def rect(slide, l, t, w, h, color, radius=False):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE, l, t, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = color; s.line.fill.background()
    return s

def txt(slide, l, t, w, h, text, size=18, color=WHITE, bold=False, align=PP_ALIGN.LEFT):
    tx = slide.shapes.add_textbox(l, t, w, h)
    tf = tx.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = text
    p.font.size = Pt(size); p.font.color.rgb = color; p.font.bold = bold
    p.font.name = "Segoe UI"; p.alignment = align
    return tx

def accent_line(slide, l, t, w, color=ACCENT):
    rect(slide, l, t, w, Inches(0.035), color)

def section_label(slide, text, y=Inches(0.8)):
    txt(slide, Inches(1.2), y, Inches(4), Inches(0.3), text, 11, DIM, True)
    accent_line(slide, Inches(1.2), y + Inches(0.35), Inches(0.5), ACCENT)

def page_num(slide, num, total=10):
    txt(slide, Inches(11.8), Inches(6.9), Inches(1), Inches(0.4),
        f"{num}/{total}", 10, DIM, align=PP_ALIGN.RIGHT)

def add_chart(slide, name, left, top, width, height=None):
    path = os.path.join(CHART_DIR, f"{name}.png")
    if height:
        slide.shapes.add_picture(path, left, top, width, height)
    else:
        slide.shapes.add_picture(path, left, top, width)


# ═════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ═════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6]); bg_fill(sl)
accent_line(sl, Inches(1.2), Inches(2.4), Inches(1.2), ACCENT)
txt(sl, Inches(1.2), Inches(2.7), Inches(10), Inches(1.2),
    "Early-Stage Prediction\nof Dry Eye Disease", 44, WHITE, True)
txt(sl, Inches(1.2), Inches(4.5), Inches(8), Inches(0.6),
    "A Machine Learning Approach Using Lifestyle, Sleep & Digital Habits", 18, MUTED)
txt(sl, Inches(1.2), Inches(6.0), Inches(4), Inches(0.4),
    "2026  ·  Research Proposal", 14, DIM)
page_num(sl, 1)


# ═════════════════════════════════════════════════════════════════════════
# SLIDE 2 — The Problem (stat cards)
# ═════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6]); bg_fill(sl)
section_label(sl, "THE PROBLEM")

stats = [
    ("344M+", "people affected\nglobally", ACCENT),
    ("75%",   "remain\nundiagnosed",       ROSE),
    ("$3.8B", "annual burden\n(US alone)",  AMBER),
    ("7+ hrs","avg daily screen\nexposure", TEAL),
]
for i, (num, desc, clr) in enumerate(stats):
    x = Inches(1.2 + i * 2.9)
    rect(sl, x, Inches(1.8), Inches(2.6), Inches(2.0), SURFACE, radius=True)
    accent_line(sl, x, Inches(1.8), Inches(2.6), clr)
    txt(sl, x + Inches(0.35), Inches(2.3), Inches(2), Inches(0.7), num, 36, clr, True)
    txt(sl, x + Inches(0.35), Inches(3.05), Inches(2), Inches(0.6), desc, 12, MUTED)

# Cost comparison chart embedded
add_chart(sl, "cost_comparison", Inches(1.2), Inches(4.3), Inches(7.5), Inches(2.8))
page_num(sl, 2)


# ═════════════════════════════════════════════════════════════════════════
# SLIDE 3 — Why Early Detection + Cost Visual
# ═════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6]); bg_fill(sl)
section_label(sl, "WHY EARLY DETECTION MATTERS")

reasons = [
    ("Irreversible Damage", "Chronic DED → corneal scarring & vision loss", ROSE),
    ("80% Preventable", "Early-stage manageable with lifestyle changes", GREEN),
    ("Gen Z at 2.5× Risk", "Digital generation: dramatically higher incidence", ACCENT),
    ("Silent Onset", "Symptoms dismissed until damage is advanced", AMBER),
]
for i, (title, desc, clr) in enumerate(reasons):
    y = Inches(1.5 + i * 1.2)
    rect(sl, Inches(1.2), y, Inches(5.5), Inches(1.0), SURFACE, radius=True)
    rect(sl, Inches(1.5), y + Inches(0.32), Inches(0.1), Inches(0.1), clr)
    txt(sl, Inches(1.8), y + Inches(0.2), Inches(4.5), Inches(0.35), title, 15, WHITE, True)
    txt(sl, Inches(1.8), y + Inches(0.55), Inches(4.5), Inches(0.4), desc, 11, MUTED)

# Symptom overlap chart on the right
add_chart(sl, "symptom_overlap", Inches(7.5), Inches(1.5), Inches(4.8), Inches(4.5))
page_num(sl, 3)


# ═════════════════════════════════════════════════════════════════════════
# SLIDE 4 — Problem Statement
# ═════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6]); bg_fill(sl)
section_label(sl, "PROBLEM STATEMENT")

rect(sl, Inches(1.2), Inches(1.8), Inches(10.9), Inches(2.2), SURFACE, radius=True)
accent_line(sl, Inches(1.2), Inches(1.8), Inches(0.06), ACCENT)
stmt = (
    '"Can we predict the onset of Dry Eye Disease using non-invasive, '
    'routinely available patient data — including sleep patterns, stress levels, '
    'digital screen habits, and ocular symptoms — enabling timely preventive '
    'intervention before irreversible damage occurs?"'
)
txt(sl, Inches(1.8), Inches(2.0), Inches(9.8), Inches(1.8), stmt, 19, WHITE)

objectives = [
    ("01", "Classification model achieving ≥ 75% F1-Score"),
    ("02", "Identify top risk factors via feature importance"),
    ("03", "Interpretable risk-scoring tool for primary care"),
    ("04", "Validate across demographics for generalizability"),
]
for i, (num, obj) in enumerate(objectives):
    x = Inches(1.2 + (i % 2) * 5.6)
    y = Inches(4.5 + (i // 2) * 0.85)
    txt(sl, x, y, Inches(0.5), Inches(0.4), num, 14, ACCENT, True)
    txt(sl, x + Inches(0.5), y, Inches(4.8), Inches(0.4), obj, 13, MUTED)

page_num(sl, 4)


# ═════════════════════════════════════════════════════════════════════════
# SLIDE 5 — Dataset with Donut Chart
# ═════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6]); bg_fill(sl)
section_label(sl, "DATASET")

# Hero stats
for i, (val, label) in enumerate([("20,000", "Records"), ("26", "Features")]):
    x = Inches(1.2 + i * 2.2)
    txt(sl, x, Inches(1.6), Inches(2), Inches(0.6), val, 38, ACCENT, True)
    txt(sl, x, Inches(2.2), Inches(2), Inches(0.3), label, 12, DIM, True)

# Donut chart for class split
add_chart(sl, "donut_class", Inches(8.5), Inches(1.0), Inches(3.5), Inches(3.5))

# Category cards
categories = [
    ("Demographics",   "Gender · Age · Height · Weight",                       ACCENT),
    ("Sleep & Stress", "Duration · Quality · Disorder · Sleepiness · Stress",  VIOLET),
    ("Cardiovascular", "Blood Pressure · Heart Rate · Steps · Activity",       GREEN),
    ("Lifestyle",      "Caffeine · Alcohol · Smoking · Medication",            AMBER),
    ("Digital",        "Screen Time · Device Before Bed · Blue-light Filter",  ROSE),
    ("Ocular",         "Eye-strain · Redness · Itchiness / Irritation",        TEAL),
]
for i, (cat, feats, clr) in enumerate(categories):
    col = i % 3
    row = i // 3
    x = Inches(1.2 + col * 3.8)
    y = Inches(3.4 + row * 1.7)
    rect(sl, x, y, Inches(3.5), Inches(1.4), SURFACE, radius=True)
    accent_line(sl, x, y, Inches(3.5), clr)
    txt(sl, x + Inches(0.3), y + Inches(0.3), Inches(3), Inches(0.35), cat, 14, WHITE, True)
    txt(sl, x + Inches(0.3), y + Inches(0.75), Inches(3), Inches(0.5), feats, 11, MUTED)

page_num(sl, 5)


# ═════════════════════════════════════════════════════════════════════════
# SLIDE 6 — Methodology Flow Diagram
# ═════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6]); bg_fill(sl)
section_label(sl, "METHODOLOGY")

# Flow diagram image — centered
add_chart(sl, "methodology_flow", Inches(0.8), Inches(1.6), Inches(11.7), Inches(2.8))

# Details below the flow
steps_detail = [
    ("01  Data Cleaning", "Handle missing values, parse BP, encode Y/N → 0/1, normalize"),
    ("02  EDA", "Distributions, correlations, class imbalance assessment (65:35)"),
    ("03  Feature Engineering", "Cramér's V (categorical), Point-Biserial r (numeric), Chi² tests"),
    ("04  Model Training", "6 classifiers, 80/20 stratified split, StandardScaler for distance-based"),
    ("05  Evaluation", "Accuracy, Precision, Recall, F1-Score, AUC-ROC, 5-fold cross-validation"),
    ("06  Deployment", "Feature importance ranking, confusion matrix analysis, risk-score tool"),
]
for i, (title, desc) in enumerate(steps_detail):
    col = i % 3
    row = i // 3
    x = Inches(1.2 + col * 3.8)
    y = Inches(4.6 + row * 1.25)
    rect(sl, x, y, Inches(3.5), Inches(1.05), SURFACE, radius=True)
    txt(sl, x + Inches(0.25), y + Inches(0.15), Inches(3.1), Inches(0.35), title, 12, ACCENT, True)
    txt(sl, x + Inches(0.25), y + Inches(0.5), Inches(3.1), Inches(0.5), desc, 10, MUTED)

page_num(sl, 6)


# ═════════════════════════════════════════════════════════════════════════
# SLIDE 7 — ML Models
# ═════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6]); bg_fill(sl)
section_label(sl, "MODELS")

models = [
    ("Logistic Regression", "Linear",         "Baseline · Interpretable coefficients",  ACCENT),
    ("Decision Tree",       "Tree-Based",     "Explainable · Visual decision paths",    GREEN),
    ("Random Forest",       "Ensemble",       "Bagging · Reduces overfitting",          TEAL),
    ("Gradient Boosting",   "Ensemble",       "Sequential error correction",            AMBER),
    ("XGBoost",             "Ensemble",       "State-of-the-art tabular performance",   ROSE),
    ("K-Nearest Neighbors", "Instance-Based", "Non-parametric · Distance-based",        VIOLET),
]
for i, (name, cat, desc, clr) in enumerate(models):
    col = i % 3
    row = i // 3
    x = Inches(1.2 + col * 3.8)
    y = Inches(1.6 + row * 2.6)
    rect(sl, x, y, Inches(3.5), Inches(2.2), SURFACE, radius=True)
    accent_line(sl, x, y, Inches(3.5), clr)
    txt(sl, x + Inches(0.35), y + Inches(0.4), Inches(3), Inches(0.4), name, 16, WHITE, True)
    rect(sl, x + Inches(0.35), y + Inches(0.95), Inches(1.4), Inches(0.32), BORDER, radius=True)
    txt(sl, x + Inches(0.5), y + Inches(0.96), Inches(1.2), Inches(0.3), cat, 10, clr, True)
    txt(sl, x + Inches(0.35), y + Inches(1.5), Inches(2.8), Inches(0.5), desc, 12, MUTED)

page_num(sl, 7)


# ═════════════════════════════════════════════════════════════════════════
# SLIDE 8 — Key Insights with Charts
# ═════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6]); bg_fill(sl)
section_label(sl, "KEY INSIGHTS — RISK FACTORS")

# Risk factors bar chart (left side)
add_chart(sl, "risk_factors", Inches(0.6), Inches(1.4), Inches(6.5), Inches(5))

# Age prevalence bar chart (right side)
add_chart(sl, "age_prevalence", Inches(7.5), Inches(1.4), Inches(5), Inches(4.2))

txt(sl, Inches(7.5), Inches(5.8), Inches(5), Inches(0.4),
    "DED Prevalence by Age Group", 12, DIM, True, PP_ALIGN.CENTER)

page_num(sl, 8)


# ═════════════════════════════════════════════════════════════════════════
# SLIDE 9 — Model Comparison Chart
# ═════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6]); bg_fill(sl)
section_label(sl, "MODEL PERFORMANCE")

# Model comparison chart — large and centered
add_chart(sl, "model_comparison", Inches(1.2), Inches(1.4), Inches(8), Inches(4.8))

# Key takeaways on the right
rect(sl, Inches(9.8), Inches(1.4), Inches(2.8), Inches(4.8), SURFACE, radius=True)
txt(sl, Inches(10.1), Inches(1.7), Inches(2.3), Inches(0.4), "Takeaways", 14, WHITE, True)

takeaways = [
    ("XGBoost", "Best overall F1 and accuracy", ROSE),
    ("Gradient\nBoosting", "Strong runner-up on all metrics", AMBER),
    ("Random\nForest", "Robust with low overfitting", TEAL),
    ("KNN", "Weakest — struggles with high dimensionality", DIM),
]
for i, (model, note, clr) in enumerate(takeaways):
    y = Inches(2.3 + i * 0.9)
    rect(sl, Inches(10.1), y, Inches(0.08), Inches(0.4), clr)
    txt(sl, Inches(10.35), y - Inches(0.02), Inches(2.2), Inches(0.3), model, 11, WHITE, True)
    txt(sl, Inches(10.35), y + Inches(0.28), Inches(2.2), Inches(0.5), note, 9, MUTED)

page_num(sl, 9)


# ═════════════════════════════════════════════════════════════════════════
# SLIDE 10 — Thank You
# ═════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6]); bg_fill(sl)
accent_line(sl, Inches(5.8), Inches(2.6), Inches(1.8), ACCENT)
txt(sl, Inches(1), Inches(2.9), Inches(11.3), Inches(1.0),
    "Thank You", 52, WHITE, True, PP_ALIGN.CENTER)
txt(sl, Inches(1), Inches(4.0), Inches(11.3), Inches(0.5),
    "Questions & Discussion", 18, MUTED, False, PP_ALIGN.CENTER)
accent_line(sl, Inches(6.1), Inches(5.0), Inches(1.2), BORDER)
txt(sl, Inches(1), Inches(5.4), Inches(11.3), Inches(0.4),
    "2026", 12, DIM, False, PP_ALIGN.CENTER)
page_num(sl, 10)


# ── Save ─────────────────────────────────────────────────────────────────
prs.save(OUT)
print(f"✅ Presentation saved: {OUT}")
print(f"   {len(prs.slides)} slides  ·  {os.path.getsize(OUT) / 1024:.0f} KB")
